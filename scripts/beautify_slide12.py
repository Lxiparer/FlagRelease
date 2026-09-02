#!/usr/bin/env python3
"""只美化第12页(index 11)：清残留孤立形状 + 4列居中拉宽 + 元数据对齐4阶段。
严格保留用户在该页手改的阶段名/术语；不触碰其它任何页。
删除旧第12页后按第5页规格重建（居中4列），新页仍落在第12位。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

F = "/home/lz/下载/FlagRelease成果汇报终版_含分支B.pptx"

NAVY=RGBColor(0x1F,0x38,0x64); BLUE=RGBColor(0x00,0x70,0xC0)
ORANGE=RGBColor(0xED,0x7D,0x31); GREEN=RGBColor(0x2E,0x9E,0x5B)
CARD=RGBColor(0xDC,0xE9,0xF7); SOFT=RGBColor(0xF7,0xF9,0xFC)
BORD=RGBColor(0xD5,0xDE,0xEA); GREY=RGBColor(0x59,0x59,0x59)
DARK=RGBColor(0x26,0x26,0x26); WHITE=RGBColor(0xFF,0xFF,0xFF)
GREENC=RGBColor(0xE0,0xF0,0xE6); FONT="微软雅黑"


def rect(sl,l,t,w,h,fill,rounded=True,line=None,lw=1.0):
    s=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                          Inches(l),Inches(t),Inches(w),Inches(h))
    if rounded:
        try: s.adjustments[0]=0.167
        except Exception: pass
    s.fill.solid(); s.fill.fore_color.rgb=fill
    if line is None: s.line.fill.background()
    else: s.line.color.rgb=line; s.line.width=Pt(lw)
    s.shadow.inherit=False
    return s


def settext(shape,runs,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,wrap=True):
    tf=shape.text_frame; tf.word_wrap=wrap; tf.vertical_anchor=anchor
    tf.margin_left=tf.margin_right=Pt(3); tf.margin_top=tf.margin_bottom=Pt(1)
    paras=runs if runs and isinstance(runs[0],list) else [runs]
    for i,para in enumerate(paras):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        for (txt,sz,b,c) in para:
            r=p.add_run(); r.text=txt; r.font.name=FONT; r.font.size=Pt(sz)
            r.font.bold=b; r.font.color.rgb=c


def textbox(sl,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,wrap=True):
    tb=sl.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    settext(tb,runs,align,anchor,wrap); return tb


def delete_slide(prs, idx):
    """删除指定索引的幻灯片（含关系）。"""
    slides = prs.slides._sldIdLst
    sldId = list(slides)[idx]
    rId = sldId.get(qn('r:id'))
    prs.part.drop_rel(rId)
    slides.remove(sldId)


def build(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    SW=13.33
    # 顶部标题区
    rect(slide,0,0,SW,1.15,WHITE,rounded=False)
    rect(slide,0.50,0.30,0.13,0.56,BLUE,rounded=False)
    textbox(slide,0.78,0.24,10.6,0.68,
            [("分支 B 工作流：三选定基线，四阶段逐版本产出",27,True,NAVY)],anchor=MSO_ANCHOR.MIDDLE)
    textbox(slide,0.80,0.84,11.0,0.30,
            [("准入镜像 gems+tree+plugin · 四阶段接力，逐阶段产出优化镜像",12,False,GREY)])
    pn=rect(slide,12.38,0.36,0.56,0.42,CARD); settext(pn,[("06",13,True,BLUE)])
    rect(slide,0,7.36,SW,0.14,NAVY,rounded=False)
    # 输入/产出 pill
    pin=rect(slide,0.70,1.32,2.20,0.50,DARK)
    settext(pin,[[("输入",11,True,WHITE),("  容器 + 模型名",10,False,CARD)]])
    pout=rect(slide,10.43,1.32,2.20,0.50,GREEN)
    settext(pout,[[("产出",11,True,WHITE),("  基线→express 镜像",10,False,GREENC)]])

    # ---- 4 阶段（严格保留用户术语）----
    stages=[
        (BLUE,  "阶段一","确定基线", ["基线镜像"],
         [("1","容器准备"),("2","环境检测"),("3","三选基线")]),
        (ORANGE,"阶段二","使能 GT", ["GT 镜像"],
         [("4","精度评测"),("5","精度调优"),("6","性能评测"),("7","性能调优"),("8","镜像发布")]),
        (GREEN, "阶段三","使能插件", ["高占比","MAX 镜像"],
         [("9","启动插件"),("10","精度评测"),("11","性能评测"),("12","镜像发布")]),
        (NAVY,  "阶段四","性能优化", ["高性能","express 镜像"],
         [("13","性能调优"),("14","精度对齐"),("15","镜像发布")]),
    ]
    n_col=len(stages); x0=0.70; gap=0.33
    col_w=(SW-2*x0-(n_col-1)*gap)/n_col   # 居中拉宽
    box_top,box_h=2.05,3.50
    xs=[x0+i*(col_w+gap) for i in range(n_col)]

    for i,(color,sn,name,caption,steps) in enumerate(stages):
        x=xs[i]
        rect(slide,x,box_top,col_w,box_h,SOFT,line=color,lw=1.3)
        head=rect(slide,x+0.12,box_top+0.12,col_w-0.24,0.66,color)
        settext(head,[[(sn,12.5,True,WHITE)],[(name,13,True,WHITE)]])
        # 产出 caption（多行术语）
        cap_para=[[("产出",10,True,color)]]+[[(t,10,True,color)] for t in caption]
        textbox(slide,x+0.12,box_top+0.80,col_w-0.24,0.34,
                cap_para,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.TOP)

        area_top=box_top+1.16; area_bot=box_top+box_h-0.12
        n=len(steps); g=0.13
        ch=min(0.80,(area_bot-area_top-(n-1)*g)/n)
        block=n*ch+(n-1)*g; start=area_top+(area_bot-area_top-block)/2
        cx=x+0.14; cw=col_w-0.28
        for j,(num,label) in enumerate(steps):
            cy=start+j*(ch+g)
            rect(slide,cx,cy,cw,ch,WHITE,line=BORD,lw=1.0)
            circ=rect(slide,cx+0.12,cy+ch/2-0.18,0.48,0.36,color)
            settext(circ,[(num,10.5,True,WHITE)])
            textbox(slide,cx+0.72,cy,cw-0.84,ch,[(label,11.5,False,DARK)],
                    align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.MIDDLE)
            if j<n-1:
                rect(slide,x+col_w/2-0.012,cy+ch,0.024,g,BORD,rounded=False)
        if i<n_col-1:
            textbox(slide,x+col_w+0.02,box_top+box_h/2-0.20,gap-0.04,0.40,
                    [("›",20,True,CARD)],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

    # 底部深蓝原则框（术语对齐：GT=V2 / MAX=V3）
    box=rect(slide,0.70,5.72,11.95,1.18,NAVY)
    settext(box,
        [[("V1 三选：",12.5,True,ORANGE),
          ("v1.1 空插件 / v1.2 厂商插件(metax) / v1.3 fl 不开算子 / none 强依赖 —— 确定性判定，冒烟测例通过即定案",11.5,False,WHITE)],
         [("子分支：",12,True,ORANGE),
          ("GT 走代码注入或 plugin 使能；MAX 走切白名单或同镜像。v1.3 场景 GT＝MAX 双 tag，阶段三整段跳过（GT 已含 plugin）",11,False,CARD)],
         [("产出门控：",12,True,ORANGE),
          ("QUALIFIED_CORE = 起服务 ∧ 精度对齐 NV（性能不阻断）；none 全失败 → 标记「模型-flagos-厂商-incompatible」",11,False,CARD)]],
        align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.MIDDLE)


def main():
    prs=Presentation(F)
    assert len(prs.slides._sldIdLst)==12, "预期12页"
    delete_slide(prs,11)          # 删旧第12页
    build(prs)                    # 重建，追加为新第12页
    prs.save(F)
    print(f"✓ 完成，共 {len(prs.slides._sldIdLst)} 页")


if __name__=="__main__":
    main()
