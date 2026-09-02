#!/usr/bin/env python3
"""向 FlagRelease 成果汇报 pptx 追加一页「分支 B 工作流」，严格复刻第5页(slide[4] 流水线页)规格。

slide[4] 设计语言（精确复刻）：
  - 顶部白标题区(h1.15) + 左侧 0070C0 竖条 + 27pt 深蓝标题 + 12pt 灰副标 + 右上 DCE9F7 圆角页码框
  - 底部深蓝细条(y7.36 h0.14)
  - 输入 pill(262626) / 产出 pill(2E9E5B)，顶行 y1.32 h0.5
  - 阶段列：F7F9FC 圆角框 + 阶段色 1.3pt 描边(top2.05 h3.5)；内含阶段色头 pill(h0.72)
  - 步骤卡：白底 D5DEEA 1pt 描边圆角 + 阶段色编号圆角标(0.55x0.38) + 11.5pt 名称
  - 底部深蓝 1F3864 圆角原则框 + 橙色引导词
  字体全微软雅黑；圆角 adj=0.167。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SRC = "/home/lz/下载/FlagRelease成果汇报终版.pptx"
OUT = "/home/lz/下载/FlagRelease成果汇报终版_含分支B.pptx"

NAVY   = RGBColor(0x1F, 0x38, 0x64)
BLUE   = RGBColor(0x00, 0x70, 0xC0)
ORANGE = RGBColor(0xED, 0x7D, 0x31)
GREEN  = RGBColor(0x2E, 0x9E, 0x5B)
PURPLE = RGBColor(0x6B, 0x39, 0xB6)
CARD   = RGBColor(0xDC, 0xE9, 0xF7)
SOFT   = RGBColor(0xF7, 0xF9, 0xFC)
BORD   = RGBColor(0xD5, 0xDE, 0xEA)
GREY   = RGBColor(0x59, 0x59, 0x59)
DARK   = RGBColor(0x26, 0x26, 0x26)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x26, 0x26, 0x26)
FONT   = "微软雅黑"


def _noshadow(s):
    s.shadow.inherit = False


def rect(slide, l, t, w, h, fill, rounded=True, line=None, line_w=1.0):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if rounded:
        try: shp.adjustments[0] = 0.167
        except Exception: pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    _noshadow(shp)
    return shp


def settext(shape, runs, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=True):
    tf = shape.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(3)
    tf.margin_top = tf.margin_bottom = Pt(1)
    paras = runs if runs and isinstance(runs[0], list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        for (text, size, bold, color) in para:
            r = p.add_run(); r.text = text
            r.font.name = FONT; r.font.size = Pt(size)
            r.font.bold = bold; r.font.color.rgb = color


def textbox(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    settext(tb, runs, align, anchor, wrap)
    return tb


def build(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    SW = 13.33

    # ---------- 顶部标题区（复刻 slide[4]）----------
    rect(slide, 0, 0, SW, 1.15, WHITE, rounded=False)
    rect(slide, 0.50, 0.30, 0.13, 0.56, BLUE, rounded=False)
    textbox(slide, 0.78, 0.24, 10.6, 0.68,
            [("分支 B 工作流：三选定基线，五版本接力产出", 27, True, NAVY)],
            anchor=MSO_ANCHOR.MIDDLE)
    textbox(slide, 0.80, 0.84, 11.0, 0.30,
            [("准入镜像 gems+tree+plugin · 五阶段接力，每阶段产出一个版本 V1–V5", 12, False, GREY)])
    pn = rect(slide, 12.38, 0.36, 0.56, 0.42, CARD)
    settext(pn, [("05", 13, True, BLUE)])
    rect(slide, 0, 7.36, SW, 0.14, NAVY, rounded=False)

    # ---------- 输入 / 产出 pill ----------
    p_in = rect(slide, 0.70, 1.32, 2.20, 0.50, INK)
    settext(p_in, [[("输入", 11, True, WHITE), ("  容器 + 模型名", 10, False, CARD)]])
    p_out = rect(slide, 10.43, 1.32, 2.20, 0.50, GREEN)
    settext(p_out, [[("产出", 11, True, WHITE), ("  V1–V5 镜像", 10, False, RGBColor(0xE0,0xF0,0xE6))]])

    # ---------- 5 阶段列 ----------
    stages = [
        (BLUE,   "阶段一", "环境就绪", "产出 V1 基础版",
         [("1", "容器准备"), ("2", "环境检测"), ("3", "三选基线")]),
        (ORANGE, "阶段二", "评测调优", "产出 V2 Pro 版",
         [("4", "精度评测"), ("5", "精度调优"), ("6", "性能评测"),
          ("7", "性能调优"), ("8", "镜像发布")]),
        (GREEN,  "阶段三", "Plugin 切换", "产出 V3 Max 版",
         [("9", "启动插件"), ("10", "精度评测"), ("11", "性能评测"), ("12", "镜像发布")]),
        (NAVY,   "阶段四", "性能优化", "产出 V4 Flag-express",
         [("13", "性能调优"), ("14", "精度对齐"), ("15", "镜像发布")]),
        (PURPLE, "阶段五", "算子扩展", "产出 V5 Royal",
         [("16", "应开尽开"), ("17", "交付发布")]),
    ]

    x0, gap = 0.70, 0.20
    col_w = (SW - 2*x0 - (len(stages)-1)*gap) / len(stages)   # ≈2.226
    box_top, box_h = 2.05, 3.50
    xs = [x0 + i*(col_w+gap) for i in range(len(stages))]

    for i, (color, sn, name, badge, steps) in enumerate(stages):
        x = xs[i]
        # 阶段外框
        rect(slide, x, box_top, col_w, box_h, SOFT, line=color, line_w=1.3)
        # 阶段色头 pill（两行：阶段N / 名称）
        head = rect(slide, x+0.12, box_top+0.12, col_w-0.24, 0.66, color)
        settext(head, [[(sn, 12.5, True, WHITE)], [(name, 13, True, WHITE)]])
        # 版本徽标 caption
        textbox(slide, x+0.12, box_top+0.82, col_w-0.24, 0.28,
                [(badge, 10, True, color)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 步骤卡区域
        area_top, area_bot = box_top+1.16, box_top+box_h-0.12
        n = len(steps)
        g = 0.13
        ch = min(0.80, (area_bot - area_top - (n-1)*g) / n)
        block = n*ch + (n-1)*g
        start = area_top + (area_bot - area_top - block)/2
        cx = x + 0.14
        cw = col_w - 0.28
        for j, (num, label) in enumerate(steps):
            cy = start + j*(ch+g)
            rect(slide, cx, cy, cw, ch, WHITE, line=BORD, line_w=1.0)
            # 编号圆角标
            circ = rect(slide, cx+0.12, cy+ch/2-0.18, 0.48, 0.36, color)
            settext(circ, [(num, 10.5, True, WHITE)])
            # 名称
            textbox(slide, cx+0.68, cy, cw-0.76, ch,
                    [(label, 11.5, False, DARK)], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
            # 列内竖向连接
            if j < n-1:
                rect(slide, x+col_w/2-0.012, cy+ch, 0.024, g, BORD, rounded=False)

        # 列间横向连接（">" 记号，居中于间隙）
        if i < len(stages)-1:
            textbox(slide, x+col_w+0.01, box_top+box_h/2-0.20, gap-0.02, 0.40,
                    [("›", 20, True, CARD)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ---------- 底部深蓝原则框（复刻 slide[4]）----------
    box = rect(slide, 0.70, 5.72, 11.95, 1.18, NAVY)
    settext(box,
        [[("V1 三选：", 12.5, True, ORANGE),
          ("v1.1 空插件 / v1.2 厂商插件(metax) / v1.3 fl 不开算子 / none 强依赖 —— 确定性判定，冒烟测例通过即定案", 11.5, False, WHITE)],
         [("子分支：", 12, True, ORANGE),
          ("V2 走代码注入或 plugin 使能；V3 走切白名单或同镜像。v1.3 场景 V2＝V3 双 tag，阶段三整段跳过（V2 已含 plugin）", 11, False, CARD)],
         [("产出门控：", 12, True, ORANGE),
          ("QUALIFIED_CORE = 起服务 ∧ 精度对齐 NV（性能不阻断）；none 全失败 → 标记「模型-flagos-厂商-incompatible」", 11, False, CARD)]],
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    return slide


def main():
    prs = Presentation(SRC)
    build(prs)
    prs.save(OUT)
    print(f"✓ 已生成: {OUT}（共 {len(prs.slides._sldIdLst)} 页）")


if __name__ == "__main__":
    main()
